import os
import tempfile
import zipfile
import streamlit as st
from pptx import Presentation
#from pptx.util import Inches, Pt
import chardet
from datetime import date  # Для определения текущей даты
#from pptx.enum.shapes import PP_PLACEHOLDER_TYPE
# Добавьте импорт для доступа к типам плейсхолдеров
#from pptx.enum.shapes import PP_PLACEHOLDER

st.set_page_config(
    page_icon="📷",
    page_title="Создание презентации из набора фотографий"
)

def create_photo_presentation(company, project_title, photo_mapping_content_bytes, photo_folder_path):
    """Creates a new presentation with slides for each photo and title."""

    i_title = 0  # Индекс названия проекта
    i_title1 = 12  # Индекс подписи левой фотографии
    i_title2 = 13 # Индекс подписи правой фотографии
    i_slide_number = 4  # Индекс номера слайда
    i_footer = 11  # Индекс футера

    if company == "эВ-групп":
        #prs_template = 'эВ_04_1.pptx' 
        prs_template = 'эВ.pptx'
    elif company == "ЭнергоCеть":
        prs_template = 'ЭС.pptx'
    else:
        raise ValueError("Для данного юридического лица нет шаблона презентации")
        

    prs = Presentation(os.path.join('template', prs_template))
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # создаём титульный слайд
    placeholders = list(slide.placeholders)
    # Обращаемся по индексу в списке (0, 1, 2...)
    #placeholders[0].text не вставляет
    placeholders[1].text = project_title
    placeholders[2].text = f"Фотоотчёт на {date.today().strftime('%d.%m.%Y')}"

    #slide_layout = prs.slide_layouts[0]  # выбираем шаблон 
    slide_layout = prs.slide_layouts[1]  # выбираем шаблон для основных слайдов презентации
    #slide = prs.slides[0]  # Первый слайд (в нашем случае единстввенный)
    #slide.placeholders[i_title].text = project_title
    
    # Detect encoding of the bytes content
    encoding = chardet.detect(photo_mapping_content_bytes)['encoding']
    # Decode the bytes using the detected encoding
    photo_mapping_content = photo_mapping_content_bytes.decode(encoding)
    lines = photo_mapping_content.strip().split('\n')

    # Координаты областей для вставки фотографий
    left1 = 517206  # левая фотография
    left2 = 6235585  # правая фотография
    top = 2200942
    width = 5442382
    height = 3996000

    N = 0 # количество вставленных в презентацию фотографий
    skipped = 0 # количество пропущенных (не добавленных в презентацию) фотографий
    for i, line in enumerate(lines):
        line = line.strip()
        if not line:
            skipped += 1
            st.warning(f"Пустая строка №{i+1}, идём дальше.")
            continue # Skip empty lines
        if ':' not in line:
            skipped += 1
            #print(f"Warning: Line '{line}' does not contain a colon, skipping.")
            st.warning(f"Строка №{i+1} '{line}' не соответствует требуемому формату [название_файла]: [описание файла], пропускаем.")
            continue # Or handle malformed lines differently


        filename, title = line.split(':', 1) # Split on first colon only
        filename = filename.strip()
        title = title.strip()
        if not title: title = filename
        
        image_path = os.path.join(photo_folder_path, filename)
        if not os.path.exists(image_path):
            skipped += 1
            #print(f"Warning: Image {image_path} not found, skipping.")
            st.warning(f"Файл {filename} не найден в ZIP-архиве, пропускаем.")
            continue # Or handle the missing image as needed

        # если не нулевая фотография и чётная, то создаём слайд
        j: int = i - skipped
        #if j>0 and j % 2 == 0:
        if j % 2 == 0:  # создаём новый слайд
            slide = prs.slides.add_slide(slide_layout)
            slide.placeholders[i_title].text = project_title 
                  
        if j % 2 == 0:
            left = left1
            slide.placeholders[i_title1].text = title
        else:
            left = left2
            slide.placeholders[i_title2].text = title

        N += 1

        # Calculate aspect ratio to fit image properly without distortion
        from PIL import Image as PILImage
        try:
            with PILImage.open(image_path) as img:
                img_width, img_height = img.size
            
            img_aspect = img_width / img_height
            shape_aspect = width / height

            if img_aspect > shape_aspect:
                # Image is wider relative to its height than the shape -> fit to width
                pic_width = width
                pic_height = int(width / img_aspect)
                pic_left = left
                pic_top = top + (height - pic_height) // 2
            else:
                # Image is taller relative to its width than the shape -> fit to height
                pic_height = height
                pic_width = int(height * img_aspect)
                pic_left = left + (width - pic_width) // 2
                pic_top = top

            slide.shapes.add_picture(image_path, pic_left, pic_top, pic_width, pic_height)
        except Exception as e:
            print(f"Error processing image {image_path}: {e}")
            # Optionally, add a text box indicating the image could not be loaded
            textbox = slide.shapes.add_textbox(left, top, width, height)
            textbox.text = f"Изображение не найдено или невозможно загрузить: {filename}"
    if N % 2 != 0: # если количество фотографий нечётное, то "обнуляем" заголовок правой части
        slide.placeholders[i_title2].text = ' '


    slide = prs.slides.add_slide(prs.slide_layouts[2])  # добавляем конечный слайд


    return prs

st.title("Создание презентации из набора фотографий")

company = st.selectbox(
    "Выберите юридическое лицо",
    ("эВ-групп", "ЭнергоCеть"),
    index=0,
    placeholder=None,
)
project_title: str = st.text_input(label="Введите название проекта (верхний заголовок на каждом слайде)", value="Мой проект")
uploaded_zip = st.file_uploader("Загрузите ZIP-архив с фотографиями", type=["zip"])
uploaded_mapping_file = st.file_uploader("НЕ ОБЯЗАТЕЛЬНО: Загрузите файл подписей к фотографиям (.txt)", type=["txt"])


if st.button("Создать презентацию"):
    if not (uploaded_zip): # Теперь проверяем только ZIP
        st.error("Загрузите ZIP-архив с фотографиями.")
    else:
        # Create a temporary directory to extract the ZIP contents
        with tempfile.TemporaryDirectory() as temp_dir:
            try:
                # Extract the uploaded ZIP file to the temporary directory
                with zipfile.ZipFile(uploaded_zip, 'r') as zip_ref:
                    zip_ref.extractall(temp_dir)
                
                # --- Генерация файла сопоставления или чтение загруженного ---
                if uploaded_mapping_file is None:
                     # Получаем отсортированный список файлов извлечённых из ZIP и сортируем
                     extracted_files = sorted([f for f in os.listdir(temp_dir) if os.path.isfile(os.path.join(temp_dir, f))])
                     
                     # Формируем строки вида "имя_файла: Добавить описание!"
                     mapping_lines = []
                     for filename in extracted_files:
                         #line = f"{filename}: ДобавитьОписание\n"
                         line = f"{filename}: {filename}\n"
                         mapping_lines.append(line)
                     
                     # Записываем строки в байтовую строку в кодировке UTF-8
                     mapping_content_str = "".join(mapping_lines)
                     mapping_content_bytes = mapping_content_str.encode('utf-8')
                     
                else: # Если файл сопоставления был загружен
                    # Читаем байты из загруженного файла
                    mapping_content_bytes = uploaded_mapping_file.read()
                
                # Вызываем функцию создания презентации с байтами содержимого файла сопоставления
                final_prs = create_photo_presentation(company, project_title, mapping_content_bytes, temp_dir)

                # Save the final presentation to a temporary file
                temp_file_path = os.path.join(tempfile.gettempdir(), "photo_report.pptx")
                final_prs.save(temp_file_path)
                st.success("Презентация создана.")
                # Provide the file for download
                with open(temp_file_path, "rb") as f:
                    st.download_button(
                        label="Скачать презентацию",
                        data=f,
                        file_name="Презентация.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                    )
                    
            except zipfile.BadZipFile:
                st.error("Загруженный файл не является действительным ZIP-архивом.")
            except Exception as e:
                st.error(f"Произошла ошибка при обработке ZIP-архива или создании презентации: {e}")

st.markdown(
        """
        <hr>
        <p style="text-align: left; color: gray;">
        <small>
        2025-2026, С.В. Медведев
        </small>
        </p>
        """,
        unsafe_allow_html=True
    )